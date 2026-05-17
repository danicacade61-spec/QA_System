"""
生成中期报告PPT
用于课程设计答辩展示
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== 颜色主题 =====
PRIMARY = RGBColor(0x1A, 0x56, 0xDB)      # 深蓝色（主色调）
PRIMARY_LIGHT = RGBColor(0x3B, 0x82, 0xF6) # 亮蓝
SECONDARY = RGBColor(0x10, 0xB9, 0x81)     # 绿色（成功/完成）
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)        # 橙色（进行中）
DARK = RGBColor(0x1E, 0x29, 0x3B)          # 深色字
GRAY = RGBColor(0x94, 0xA3, 0xB8)          # 灰色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)     # 浅背景
CARD_BG = RGBColor(0xF0, 0xF4, 0xF8)      # 卡片背景


def add_bg_shape(slide, color=PRIMARY):
    """添加背景色块"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    """添加形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, 
                 bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=13, 
                   color=DARK, spacing=Pt(6), font_name="微软雅黑"):
    """添加多行文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # 支持区分标题和正文
        if isinstance(line, tuple):
            text, bold, size, clr = line
            p.text = text
            p.font.bold = bold
            p.font.size = Pt(size)
            p.font.color.rgb = clr
        else:
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
        
        p.font.name = font_name
        p.space_after = spacing
    
    return txBox


def add_card(slide, left, top, width, height, title, content_lines, 
             card_color=PRIMARY, title_color=WHITE):
    """添加卡片式内容块"""
    # 卡片背景
    card = add_shape(slide, left, top, width, height, CARD_BG)
    card.shadow.inherit = False
    
    # 顶部色条
    add_shape(slide, left, top, width, Cm(0.5), card_color)
    
    # 标题
    add_text_box(slide, left + Cm(0.5), top + Cm(0.3), width - Cm(1), Cm(0.8),
                 title, font_size=13, bold=True, color=DARK)
    
    # 内容
    y_offset = top + Cm(1.3)
    for line in content_lines:
        add_text_box(slide, left + Cm(0.5), y_offset, width - Cm(1), Cm(0.5),
                     f"• {line}", font_size=11, color=DARK)
        y_offset += Cm(0.5)


def create_slide_title(slide, title, subtitle=None):
    """创建幻灯片标题区域"""
    # 顶部色带
    add_shape(slide, Cm(0), Cm(0), Cm(25.4), Cm(2.2), PRIMARY)
    
    # 标题文字
    add_text_box(slide, Cm(1.5), Cm(0.3), Cm(22), Cm(1.2),
                 title, font_size=26, bold=True, color=WHITE)
    
    if subtitle:
        add_text_box(slide, Cm(1.5), Cm(1.3), Cm(22), Cm(0.8),
                     subtitle, font_size=13, color=RGBColor(0xBF, 0xDB, 0xFE))


def create_ppt():
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(14.29)
    
    # ======================== Slide 1: 封面 ========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    add_bg_shape(slide, PRIMARY)
    
    # 装饰线
    add_shape(slide, Cm(2), Cm(3), Cm(21.4), Cm(0.1), SECONDARY)
    add_shape(slide, Cm(2), Cm(10.5), Cm(21.4), Cm(0.1), SECONDARY)
    
    # 标题
    add_text_box(slide, Cm(2), Cm(3.5), Cm(21.4), Cm(1.5),
                 "基于大语言模型的智能问答方法研究与实现",
                 font_size=30, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # 副标题
    add_text_box(slide, Cm(2), Cm(5.5), Cm(21.4), Cm(1),
                 "课程设计中期报告",
                 font_size=20, color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)
    
    # 学校信息
    add_text_box(slide, Cm(2), Cm(7.5), Cm(21.4), Cm(0.6),
                 "西南交通大学  计算机与人工智能学院",
                 font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # 底部信息
    add_text_box(slide, Cm(2), Cm(11), Cm(21.4), Cm(0.5),
                 "学号：2025212413 | 姓名：毛树林 | 课程：自然语言处理",
                 font_size=13, color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Cm(2), Cm(11.6), Cm(21.4), Cm(0.5),
                 "二零二六年五月",
                 font_size=13, color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)
    
    # ======================== Slide 2: 目录 ========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_title(slide, "目  录", "CONTENTS")
    
    items = [
        ("01", "项目背景与目标", "研究动机、核心问题"),
        ("02", "技术原理概述", "大语言模型、RAG架构"),
        ("03", "系统架构设计", "模块划分、技术栈"),
        ("04", "核心模块实现", "各模块关键实现细节"),
        ("05", "前后端实现", "API服务、Web界面"),
        ("06", "进度完成情况", "已完成功能、当前进度"),
        ("07", "当前问题与后续计划", "问题分析与后续安排"),
    ]
    
    for i, (num, title, desc) in enumerate(items):
        y = Cm(3.2) + Cm(1.4) * i
        
        # 序号圆
        circle = add_shape(slide, Cm(2), y, Cm(1), Cm(1), PRIMARY, MSO_SHAPE.OVAL)
        add_text_box(slide, Cm(2), y + Cm(0.1), Cm(1), Cm(0.8),
                     num, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        
        # 标题
        add_text_box(slide, Cm(3.5), y, Cm(8), Cm(0.6),
                     title, font_size=16, bold=True, color=DARK)
        
        # 描述
        add_text_box(slide, Cm(3.5), y + Cm(0.6), Cm(8), Cm(0.4),
                     desc, font_size=11, color=GRAY)
    
    # ======================== Slide 3: 项目背景与目标 ========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_title(slide, "项目背景与目标", "BACKGROUND & OBJECTIVES")
    
    # 左侧：背景
    add_card(slide, Cm(1.5), Cm(3), Cm(11), Cm(5.5),
             "🎯 研究背景", [
                 "大语言模型（LLM）在自然语言处理领域取得突破",
                 "智能问答系统是LLM最具应用价值的场景之一",
                 "传统问答系统依赖规则/模板，泛化能力有限",
                 "RAG（检索增强生成）成为行业主流方案",
             ])
    
    # 右侧：核心目标
    add_card(slide, Cm(13), Cm(3), Cm(11), Cm(5.5),
             "⚡ 核心目标", [
                 "研究LLM实现智能问答的技术原理",
                 "实现一个完整的智能问答系统原型",
                 "支持文档知识库的检索增强问答",
                 "提供Web界面与RESTful API交互方式",
             ])
    
    # 底部：两个核心问题
    add_card(slide, Cm(1.5), Cm(9), Cm(22.5), Cm(4),
             "❓ 研究的两个核心问题", [
                 "原理层面：Transformer架构如何处理自然语言？预训练如何赋予模型问答能力？自回归生成如何逐步产出回答？",
                 "工程层面：如何将LLM落地为可用系统？包括模型部署、后端推理服务、前端交互界面的完整集成",
             ], card_color=ACCENT)
    
    # ======================== Slide 4: 大语言模型技术原理 ========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_title(slide, "大语言模型技术原理", "LLM TECHNICAL PRINCIPLES")
    
    # 三个阶段示意
    stages = [
        ("预训练", "海量文本语料\n自监督学习\n（预测下一个Token）\n→ 获得语言知识"),
        ("微调", "指令数据\n监督微调\n（对话/问答格式）\n→ 对齐人类意图"),
        ("推理生成", "给定输入\n自回归解码\n（逐Token生成）\n→ 产出回答"),
    ]
    
    for i, (title, desc) in enumerate(stages):
        x = Cm(1.5) + Cm(7.8) * i
        # 阶段框
        box = add_shape(slide, x, Cm(3.2), Cm(7), Cm(5), CARD_BG)
        
        # 顶部色条
        colors = [PRIMARY, SECONDARY, ACCENT]
        add_shape(slide, x, Cm(3.2), Cm(7), Cm(0.5), colors[i])
        
        # 阶段标题
        add_text_box(slide, x + Cm(0.5), Cm(3.8), Cm(6), Cm(0.6),
                     f"阶段{i+1}：{title}", font_size=15, bold=True, color=colors[i])
        
        # 描述
        add_text_box(slide, x + Cm(0.5), Cm(4.8), Cm(6), Cm(3),
                     desc, font_size=11, color=DARK)
        
        # 箭头（除了最后一个）
        if i < 2:
            add_text_box(slide, x + Cm(7.2), Cm(5), Cm(0.6), Cm(0.5),
                         "→", font_size=24, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    
    # 底部：关键技术
    add_card(slide, Cm(1.5), Cm(8.8), Cm(22.5), Cm(4),
             "🔧 关键技术", [
                 "Transformer架构：自注意力机制（Self-Attention）+ 前馈神经网络（FFN），并行处理序列",
                 "自回归生成：逐个预测下一个Token，使用KV-Cache加速推理（每次生成一个token）",
                 "RAG架构：检索 + 生成结合，在不重新训练模型的情况下引入外部知识",
             ], card_color=PRIMARY_LIGHT)
    
    # ======================== Slide 5: 系统架构设计 ========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_title(slide, "系统架构设计", "SYSTEM ARCHITECTURE")
    
    # 架构流程图 - 使用形状和箭头
    steps = [
        ("输入层", "用户提问\n(Web/API)"),
        ("检索层", "向量相似度\n(FAISS)"),
        ("生成层", "LLM推理\n(智谱/OpenAI)"),
        ("输出层", "回答返回\n(流式/非流式)"),
    ]
    
    for i, (title, desc) in enumerate(steps):
        x = Cm(1.2) + Cm(5.8) * i
        # 模块框
        box = add_shape(slide, x, Cm(3.3), Cm(5), Cm(4), WHITE)
        box.line.color.rgb = PRIMARY
        box.line.width = Pt(2)
        
        # 顶部标题
        add_shape(slide, x, Cm(3.3), Cm(5), Cm(0.8), PRIMARY)
        add_text_box(slide, x, Cm(3.3), Cm(5), Cm(0.8),
                     title, font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        
        # 描述
        add_text_box(slide, x + Cm(0.3), Cm(4.5), Cm(4.4), Cm(2.5),
                     desc, font_size=11, color=DARK, alignment=PP_ALIGN.CENTER)
        
        # 箭头
        if i < 3:
            add_text_box(slide, x + Cm(5.2), Cm(4.5), Cm(0.6), Cm(0.5),
                         "▶", font_size=18, color=SECONDARY, alignment=PP_ALIGN.CENTER)
    
    # 分层架构
    add_text_box(slide, Cm(1.5), Cm(7.8), Cm(22), Cm(0.5),
                 "▎分层模块架构", font_size=14, bold=True, color=DARK)
    
    layers = [
        ("数据层", "文档加载（TXT/PDF/DOCX/MD）→ 文本分块 → 向量化 → 索引存储"),
        ("检索层", "FAISS向量检索 → 相关性排序 → 阈值过滤 → 上下文召回"),
        ("生成层", "多LLM后端支持 → 上下文构建 → Prompt工程 → 流式生成"),
        ("服务层", "FastAPI RESTful API → 知识库管理 → 对话历史管理 → CORS"),
        ("展示层", "Streamlit交互界面：对话 / 知识库管理 / 检索预览 三标签页"),
    ]
    
    for i, (layer, desc) in enumerate(layers):
        y = Cm(8.5) + Cm(1.05) * i
        colors = [PRIMARY, RGBColor(0x25, 0x6E, 0xEB), SECONDARY, ACCENT, RGBColor(0x8B, 0x5C, 0xF6)]
        
        # 标签
        tag = add_shape(slide, Cm(1.5), y, Cm(2.8), Cm(0.8), colors[i])
        add_text_box(slide, Cm(1.5), y + Cm(0.1), Cm(2.8), Cm(0.6),
                     layer, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        
        # 描述
        add_text_box(slide, Cm(4.8), y + Cm(0.1), Cm(19), Cm(0.6),
                     desc, font_size=10, color=DARK)
    
    # ======================== Slide 6: 核心模块实现 ========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_title(slide, "核心模块实现", "CORE MODULE IMPLEMENTATION")
    
    modules = [
        ("📄 文档处理", [
            "支持TXT/PDF/DOCX/MD四种格式",
            "段落级+句子级混合分块策略",
            "默认块大小256字符，重叠32字符",
            "自动过滤过短块（<10字符）",
        ]),
        ("🔤 嵌入向量", [
            "BAAI/bge-small-zh-v1.5模型",
            "输出512维向量，L2归一化",
            "sentence-transformers批量编码",
            "延迟加载策略节省启动资源",
        ]),
        ("🗄️ 向量存储", [
            "FAISS + IndexFlatIP索引",
            "相似度阈值过滤（默认0.3）",
            "持久化到磁盘，自动恢复",
            "支持增删改查操作",
        ]),
    ]
    
    for i, (title, items) in enumerate(modules):
        x = Cm(1.2) + Cm(7.9) * i
        add_card(slide, x, Cm(3), Cm(7.3), Cm(4.5), title, items)
    
    # LLM客户端 + RAG引擎卡片 - 下方
    more_modules = [
        ("🤖 LLM客户端", [
            "策略模式 + 工厂方法设计",
            "三种后端：OpenAI / 智谱 / 本地",
            "统一抽象接口，支持流式输出",
            "当前主力：智谱GLM-5.1 API",
        ]),
        ("🔗 RAG引擎", [
            "完整的检索增强生成流程",
            "检索→上下文构建→Prompt组装→生成",
            "保留最近3轮对话历史",
            "流式/非流式双接口支持",
        ]),
    ]
    
    for i, (title, items) in enumerate(more_modules):
        x = Cm(1.2) + Cm(12.3) * i
        add_card(slide, x, Cm(8), Cm(11.7), Cm(5.2), title, items, card_color=PRIMARY_LIGHT)
    
    # ======================== Slide 7: 前后端实现 ========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_title(slide, "前后端实现", "FRONTEND & BACKEND")
    
    # 左侧：后端API
    add_card(slide, Cm(1.5), Cm(3), Cm(11), Cm(5),
             "⚙️ FastAPI 后端服务", [
                 "RESTful API，部署于8000端口",
                 "问答查询（含SSE流式响应）",
                 "知识库管理（上传/删除/统计/清空）",
                 "自动生成Swagger交互式文档",
                 "CORS跨域支持 + Pydantic数据验证",
             ])
    
    # 右侧：前端界面
    add_card(slide, Cm(13), Cm(3), Cm(11), Cm(5),
             "🎨 Streamlit 前端界面", [
                 "部署于8501端口，三个功能标签页",
                 "对话页：多轮对话 + RAG检索开关",
                 "知识库管理：文件上传/批量导入/查看",
                 "检索预览：直观查看向量检索结果",
                 "侧边栏：系统状态/参数/控制面板",
             ])
    
    # 底部：评估模块
    add_card(slide, Cm(1.5), Cm(8.5), Cm(22.5), Cm(4.5),
             "📊 评估模块", [
                 "测试数据加载：支持JSON/JSONL/TXT格式",
                 "自动评估指标：Exact Match / Token F1 / ROUGE-L（最长公共子序列）",
                 "性能记录：响应时间、检索覆盖率",
                 "结果输出：自动保存JSON和TXT格式报告，支持批量评估与汇总统计",
             ], card_color=SECONDARY)
    
    # ======================== Slide 8: 进度完成情况 ========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_title(slide, "进度完成情况", "PROGRESS STATUS")
    
    # 三个阶段的进度条
    phases = [
        ("第一阶段：文献调研与技术方案", "✅ 已完成 100%", [
            "✓ 理解Transformer、预训练、自回归生成原理",
            "✓ 确定RAG架构作为技术方案",
            "✓ 选定智谱API和BGE嵌入模型",
        ], SECONDARY),
        ("第二阶段：系统设计与开发", "✅ 已完成 ~95%", [
            "✓ 全部核心模块编码完成",
            "✓ 文档处理/嵌入/向量/LLM/RAG/API/界面/评估",
            "⟳ 知识库文档加载待完成",
        ], PRIMARY),
        ("第三阶段：测试与评估", "⟳ 进行中 ~30%", [
            "✓ 评估模块开发完成",
            "✓ 测试数据集已创建（10个样本）",
            "⟳ 系统性测试评估进行中",
        ], ACCENT),
    ]
    
    for i, (title, status, items, color) in enumerate(phases):
        y = Cm(3) + Cm(3.5) * i
        
        # 阶段卡片
        card = add_shape(slide, Cm(1.5), y, Cm(22.5), Cm(3.2), CARD_BG)
        
        # 左侧色条
        add_shape(slide, Cm(1.5), y, Cm(0.3), Cm(3.2), color)
        
        # 阶段标题
        add_text_box(slide, Cm(2.5), y + Cm(0.2), Cm(12), Cm(0.5),
                     title, font_size=14, bold=True, color=color)
        
        # 进度状态
        add_text_box(slide, Cm(15), y + Cm(0.2), Cm(8), Cm(0.5),
                     status, font_size=12, bold=True, color=color, alignment=PP_ALIGN.RIGHT)
        
        # 条目
        item_text = "  ".join(items)
        add_text_box(slide, Cm(2.5), y + Cm(0.9), Cm(20.5), Cm(2),
                     item_text, font_size=10, color=DARK)
    
    # ======================== Slide 9: 功能清单 ========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_title(slide, "功能实现清单", "FEATURE COMPLETION LIST")
    
    features = [
        ("文档处理", "TXT/PDF/DOCX/MD加载与分块", "✅"),
        ("嵌入向量", "BGE-small-zh + FAISS索引", "✅"),
        ("LLM客户端", "OpenAI/智谱/本地三种后端", "✅"),
        ("RAG引擎", "检索生成 + 流式 + 多轮对话", "✅"),
        ("API服务", "FastAPI + SSE流式响应", "✅"),
        ("前端界面", "Streamlit三标签页", "✅"),
        ("评估模块", "EM / F1 / ROUGE-L", "✅"),
        ("知识库管理", "上传/删除/清空/统计", "✅"),
        ("系统配置", "多提供者/参数可配置", "✅"),
    ]
    
    # 表头
    headers = ["功能模块", "功能描述", "状态"]
    header_colors = [PRIMARY, PRIMARY, PRIMARY]
    col_widths = [Cm(5), Cm(14), Cm(3)]
    
    for j, (header, h_color, cw) in enumerate(zip(headers, header_colors, col_widths)):
        x = Cm(1.5) + sum(col_widths[:j])
        hdr = add_shape(slide, x, Cm(3.2), cw, Cm(0.8), h_color)
        add_text_box(slide, x, Cm(3.2), cw, Cm(0.8),
                     header, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # 表格行
    for i, (module, desc, status) in enumerate(features):
        y = Cm(4.2) + Cm(0.9) * i
        bg = WHITE if i % 2 == 0 else CARD_BG
        
        for j, (text, cw) in enumerate(zip([module, desc, status], col_widths)):
            x = Cm(1.5) + sum(col_widths[:j])
            cell = add_shape(slide, x, y, cw, Cm(0.8), bg)
            cell.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            cell.line.width = Pt(0.5)
            clr = SECONDARY if text == "✅" else DARK
            bld = text == "✅"
            add_text_box(slide, x, y + Cm(0.1), cw, Cm(0.6),
                         text, font_size=11, bold=bld, color=clr, alignment=PP_ALIGN.CENTER)
    
    # 总结
    add_text_box(slide, Cm(1.5), Cm(12), Cm(22.5), Cm(1),
                 "📌 全部9大功能模块均已完成开发，系统主体框架已完备，进入测试评估与报告撰写阶段",
                 font_size=13, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    
    # ======================== Slide 10: 当前问题 ========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_title(slide, "当前问题与解决方案", "CURRENT ISSUES & SOLUTIONS")
    
    issues = [
        ("问题1", "知识库为空，未加载文档", "RAG检索无法验证", "后续上传课程资料构建知识库"),
        ("问题2", "测试数据仅10条", "评估统计意义有限", "扩展至100-200条测试数据"),
        ("问题3", "本地模型推理慢", "影响体验", "以API为主，模型量化加速"),
        ("问题4", "尚未人工评估", "缺人工维度验证", "从准确/相关/流畅三维度评分"),
        ("问题5", "分块策略待优化", "影响检索准确率", "对比不同分块大小和重叠策略"),
    ]
    
    # 表头
    issue_headers = ["序号", "问题", "影响", "方案"]
    issue_widths = [Cm(2.5), Cm(7.5), Cm(6), Cm(6.5)]
    
    for j, (h, cw) in enumerate(zip(issue_headers, issue_widths)):
        x = Cm(1.2) + sum(issue_widths[:j])
        hdr = add_shape(slide, x, Cm(3.2), cw, Cm(0.8), PRIMARY)
        add_text_box(slide, x, Cm(3.2), cw, Cm(0.8),
                     h, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    for i, (num, problem, impact, solution) in enumerate(issues):
        y = Cm(4.2) + Cm(1.6) * i
        bg = WHITE if i % 2 == 0 else CARD_BG
        
        for j, (text, cw) in enumerate(zip([num, problem, impact, solution], issue_widths)):
            x = Cm(1.2) + sum(issue_widths[:j])
            cell = add_shape(slide, x, y, cw, Cm(1.4), bg)
            cell.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            cell.line.width = Pt(0.5)
            fs = 11 if j != 3 else 10
            add_text_box(slide, x + Cm(0.2), y + Cm(0.2), cw - Cm(0.4), Cm(1),
                         text, font_size=fs, color=DARK, alignment=PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT)
    
    # ======================== Slide 11: 后续计划 ========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_title(slide, "后续工作计划", "FUTURE WORK PLAN")
    
    schedule = [
        ("第11周", "本周", "📥 构建知识库\n📝 扩展测试数据", SECONDARY),
        ("第12周", "下周", "🧪 系统性测试评估\n🔧 参数调优", PRIMARY),
        ("第13周", "第3周", "👥 人工评估\n🐛 Bug修复与优化", ACCENT),
        ("第14~15周", "第4-5周", "✍️ 撰写报告初稿", RGBColor(0x8B, 0x5C, 0xF6)),
        ("第16~17周", "第6-7周", "📋 报告修改完善\n🎤 准备答辩", RGBColor(0xEC, 0x48, 0x99)),
    ]
    
    for i, (time, label, tasks, color) in enumerate(schedule):
        x = Cm(1.2) + Cm(4.8) * i
        
        # 时间卡片
        card = add_shape(slide, x, Cm(3.2), Cm(4.3), Cm(7), CARD_BG)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        # 周次标题
        add_shape(slide, x, Cm(3.2), Cm(4.3), Cm(1.2), color)
        add_text_box(slide, x, Cm(3.2), Cm(4.3), Cm(0.6),
                     time, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Cm(3.8), Cm(4.3), Cm(0.5),
                     label, font_size=9, color=WHITE, alignment=PP_ALIGN.CENTER)
        
        # 任务
        add_text_box(slide, x + Cm(0.3), Cm(4.8), Cm(3.7), Cm(5),
                     tasks, font_size=11, color=DARK)
    
    # 预期成果
    add_card(slide, Cm(1.2), Cm(10.5), Cm(23), Cm(3),
             "🏆 预期最终成果", [
                 "一个可运行的智能问答系统原型（Web + API双交互方式）",
                 "系统性测试评估报告（自动指标 + 人工评分）",
                 "完整的课程设计报告（原理阐述 + 设计方案 + 实现细节 + 测试结果）",
             ], card_color=SECONDARY)
    
    # ======================== Slide 12: 技术栈总览 ========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_title(slide, "技术栈总览", "TECHNOLOGY STACK")
    
    techs = [
        ("编程语言", "Python 3.10+"),
        ("Web框架", "FastAPI + Uvicorn"),
        ("前端框架", "Streamlit"),
        ("向量数据库", "FAISS (IndexFlatIP)"),
        ("嵌入模型", "BAAI/bge-small-zh-v1.5"),
        ("LLM提供者", "智谱GLM-5.1 / OpenAI / 本地Qwen2"),
        ("文档解析", "pypdf / python-docx"),
        ("评估工具", "rouge-score / 自建F1计算"),
        ("开发工具", "Git / VS Code / Conda"),
    ]
    
    col1 = techs[:5]
    col2 = techs[5:]
    
    for col_idx, col in enumerate([col1, col2]):
        for i, (name, detail) in enumerate(col):
            y = Cm(3.2) + Cm(1.6) * i
            x = Cm(1.5) + Cm(12) * col_idx
            
            # 标签
            tag = add_shape(slide, x, y, Cm(3.5), Cm(1.2), PRIMARY)
            add_text_box(slide, x, y + Cm(0.2), Cm(3.5), Cm(0.8),
                         name, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
            
            # 值
            val = add_shape(slide, x + Cm(3.5), y, Cm(7.5), Cm(1.2), CARD_BG)
            add_text_box(slide, x + Cm(3.8), y + Cm(0.2), Cm(7), Cm(0.8),
                         detail, font_size=11, color=DARK)
    
    # ======================== Slide 13: 谢谢 ========================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide, PRIMARY)
    
    add_shape(slide, Cm(2), Cm(4), Cm(21.4), Cm(0.1), SECONDARY)
    add_shape(slide, Cm(2), Cm(10), Cm(21.4), Cm(0.1), SECONDARY)
    
    add_text_box(slide, Cm(2), Cm(5), Cm(21.4), Cm(1.5),
                 "感谢聆听", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Cm(2), Cm(6.8), Cm(21.4), Cm(0.8),
                 "THANK YOU", font_size=18, color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Cm(2), Cm(8), Cm(21.4), Cm(1),
                 "欢迎提问与交流", font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Cm(2), Cm(10.5), Cm(21.4), Cm(0.5),
                 "学号：2025212413 | 姓名：毛树林 | 课程：自然语言处理",
                 font_size=13, color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)
    
    # ===== 保存 =====
    output_path = os.path.join(os.path.dirname(__file__), "中期报告_基于大语言模型的智能问答方法研究与实现.pptx")
    prs.save(output_path)
    print(f"✅ PPT已生成: {output_path}")


if __name__ == "__main__":
    create_ppt()
